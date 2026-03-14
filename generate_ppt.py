from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# --- THEME COLORS ---
COLOR_BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
COLOR_ACCENT     = RGBColor(0x00, 0xC2, 0xFF)   # electric cyan
COLOR_ACCENT2    = RGBColor(0x7B, 0x2F, 0xFF)   # violet
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GREY = RGBColor(0xB0, 0xBE, 0xCC)
COLOR_CARD       = RGBColor(0x14, 0x2A, 0x40)   # slightly lighter navy for cards

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # completely blank

# ─── HELPERS ────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=COLOR_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, left, top, width, height,
                  font_size=14, color=COLOR_WHITE, leading_color=COLOR_ACCENT,
                  leading_char="▸ ", bold_first=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        # leading bullet
        r1 = p.add_run()
        r1.text = leading_char
        r1.font.size = Pt(font_size)
        r1.font.color.rgb = leading_color
        r1.font.bold = True
        # content
        r2 = p.add_run()
        r2.text = line
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color
        r2.font.bold = bold_first and i == 0

def add_divider(slide, top, color=COLOR_ACCENT):
    bar = add_rect(slide, 0.6, top, 12.1, 0.04, fill_color=color)
    return bar

def tag_pill(slide, text, left, top, width=1.8, height=0.38, bg=COLOR_ACCENT2):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_text(slide, text, left, top + 0.03, width, height - 0.06,
             font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# ─── SLIDE BUILDERS ─────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)

    # gradient-like left bar
    add_rect(slide, 0, 0, 0.45, 7.5, fill_color=COLOR_ACCENT2)
    add_rect(slide, 0.45, 0, 0.06, 7.5, fill_color=COLOR_ACCENT)

    # decorative corner circle (fake with a rect/rounded look)
    add_rect(slide, 9.8, 5.2, 3.5, 3.5, fill_color=RGBColor(0x00, 0x80, 0xAA))

    add_text(slide, "NeoConnect", 1.2, 1.4, 10, 1.2,
             font_size=54, bold=True, color=COLOR_WHITE)
    add_text(slide, "Staff Feedback & Transparency Platform", 1.2, 2.65, 10, 0.6,
             font_size=22, color=COLOR_ACCENT)
    add_divider(slide, 3.4)
    add_text(slide, "Development Approach  |  Architecture  |  Challenges  |  Roadmap",
             1.2, 3.55, 11, 0.55, font_size=14, color=COLOR_LIGHT_GREY, italic=True)

    tag_pill(slide, "Full-Stack", 1.2, 4.4)
    tag_pill(slide, "Next.js + Node.js", 3.15, 4.4, width=2.3)
    tag_pill(slide, "MongoDB Atlas", 5.6, 4.4, width=2.1)

    add_text(slide, "Varshith Gorrepati  •  March 2026", 1.2, 6.6, 10, 0.5,
             font_size=11, color=COLOR_LIGHT_GREY, italic=True)


def slide_problem(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=COLOR_CARD)
    add_text(slide, "The Problem", 0.6, 0.15, 8, 0.65,
             font_size=28, bold=True, color=COLOR_ACCENT)
    add_divider(slide, 1.05)

    problems = [
        "Staff complaints get lost in email threads with no accountability.",
        "No structured system for tracking case resolution timelines.",
        "Management lacks data to identify recurring departmental issues.",
        "Anonymous feedback channels are either absent or untrusted.",
        "Meeting minutes and decisions are not easily accessible to all staff.",
    ]
    add_multiline(slide, problems, 0.6, 1.3, 12.1, 5.5, font_size=17,
                  leading_color=RGBColor(0xFF, 0x5C, 0x5C), leading_char="✖  ")

    add_text(slide, "Without a centralised platform, workplace trust erodes silently.",
             0.6, 6.5, 12, 0.55, font_size=13, color=COLOR_ACCENT, italic=True)


def slide_solution(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=COLOR_CARD)
    add_text(slide, "The Solution — NeoConnect", 0.6, 0.15, 12, 0.65,
             font_size=28, bold=True, color=COLOR_ACCENT)
    add_divider(slide, 1.05)

    # 4 solution cards
    cards = [
        ("🔐", "Secure\nRole-Based Login",   "Staff, Manager,\nSecretariat & Admin"),
        ("📋", "Structured\nCase Submission", "Anonymous toggle,\nfile uploads, tracking"),
        ("📊", "Analytics\nDashboard",        "Hotspot detection\n& trend analysis"),
        ("🔔", "Auto\nEscalation",            "7-day cron job\nautomation"),
    ]
    for i, (icon, title, desc) in enumerate(cards):
        x = 0.4 + i * 3.2
        add_rect(slide, x, 1.3, 2.95, 4.5, fill_color=COLOR_CARD,
                 line_color=COLOR_ACCENT, line_width=Pt(1.5))
        add_text(slide, icon,  x+0.9, 1.5,  2.2, 0.8, font_size=30, align=PP_ALIGN.CENTER)
        add_text(slide, title, x+0.1, 2.4,  2.75, 0.9, font_size=15, bold=True,
                 color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc,  x+0.1, 3.4,  2.75, 1.9, font_size=13,
                 color=COLOR_LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_text(slide, "One platform — four roles — zero missed complaints.",
             0.6, 6.35, 12, 0.55, font_size=13, color=COLOR_ACCENT2, italic=True)


def slide_approach(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=COLOR_CARD)
    add_text(slide, "Development Approach", 0.6, 0.15, 12, 0.65,
             font_size=28, bold=True, color=COLOR_ACCENT)
    add_divider(slide, 1.05)

    steps = [
        ("1. API-First Design",      "All REST endpoints (Auth, Cases, Analytics, Polls) were defined before the UI, ensuring clean data contracts."),
        ("2. Role Modelling",        "Identified 4 user roles with differing permissions early to design a scalable RBAC middleware layer."),
        ("3. Decoupled Frontend",    "Next.js App Router kept fully independent from Express API, allowing the frontend to be swapped without touching the backend."),
        ("4. Schema-Driven DB",      "Mongoose schemas were finalised before feature coding to prevent data inconsistencies across modules."),
        ("5. Iterative UI Polish",   "Each page was progressively enhanced with Recharts, Tailwind animations, and glass-morphism card design."),
    ]
    for i, (heading, detail) in enumerate(steps):
        y = 1.25 + i * 1.05
        add_rect(slide, 0.45, y, 0.08, 0.75, fill_color=COLOR_ACCENT)
        add_text(slide, heading, 0.7, y,       5.5, 0.4, font_size=14, bold=True, color=COLOR_ACCENT)
        add_text(slide, detail,  0.7, y + 0.38, 12, 0.55, font_size=12, color=COLOR_LIGHT_GREY)


def slide_techstack(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=COLOR_CARD)
    add_text(slide, "Technical Architecture", 0.6, 0.15, 12, 0.65,
             font_size=28, bold=True, color=COLOR_ACCENT)
    add_divider(slide, 1.05)

    layers = [
        ("FRONTEND",  "Next.js 14 (App Router) · React · Tailwind CSS · Recharts",              COLOR_ACCENT),
        ("API LAYER", "Node.js · Express.js · Express-Validator · Morgan · Multer",              COLOR_ACCENT2),
        ("AUTH",      "JWT (jsonwebtoken) · Bcrypt · Role-Based Middleware Guards",              RGBColor(0xFF, 0xA5, 0x00)),
        ("DATABASE",  "MongoDB Atlas · Mongoose ODM · Cloud replication & backups",             RGBColor(0x00, 0xC8, 0x5A)),
        ("INFRA",     "GitHub → Render (Backend CI/CD) · Vercel (Frontend) · Cron Jobs",       RGBColor(0xFF, 0x5C, 0x8D)),
    ]
    for i, (label, tech, color) in enumerate(layers):
        y = 1.25 + i * 1.05
        add_rect(slide, 0.45, y, 1.5, 0.68, fill_color=color)
        add_text(slide, label, 0.45, y+0.14, 1.5, 0.4,
                 font_size=11, bold=True, color=COLOR_BG_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, tech, 2.1, y+0.1, 10.6, 0.55,
                 font_size=13, color=COLOR_WHITE)


def slide_features(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=COLOR_CARD)
    add_text(slide, "Core Features by Role", 0.6, 0.15, 12, 0.65,
             font_size=28, bold=True, color=COLOR_ACCENT)
    add_divider(slide, 1.05)

    roles = [
        ("👤 Staff",             ["Submit anonymous complaints", "Attach PDF/image evidence", "Track case status", "Vote in polls"]),
        ("🗂 Case Manager",      ["View assigned cases", "Update status & add notes", "Respond within SLA window"]),
        ("🏛 Secretariat/Mgmt", ["View & filter all cases", "Assign to case managers", "Upload meeting minutes", "Create polls"]),
        ("⚙️ Admin",            ["Manage users & roles", "Add departments", "Full system access"]),
    ]
    for i, (role, items) in enumerate(roles):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.5
        y = 1.25 + row * 2.9
        add_rect(slide, x, y, 6.15, 2.6, fill_color=COLOR_CARD,
                 line_color=COLOR_ACCENT2, line_width=Pt(1.2))
        add_text(slide, role, x+0.15, y+0.1, 5.8, 0.45,
                 font_size=15, bold=True, color=COLOR_ACCENT)
        add_multiline(slide, items, x+0.15, y+0.55, 5.8, 1.9,
                      font_size=12, leading_char="• ", leading_color=COLOR_ACCENT2)


def slide_challenges(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=COLOR_CARD)
    add_text(slide, "Challenges & Solutions", 0.6, 0.15, 12, 0.65,
             font_size=28, bold=True, color=COLOR_ACCENT)
    add_divider(slide, 1.05)

    items = [
        ("RBAC Complexity",
         "4 user roles with overlapping permissions required custom JWT middleware that decodes role claims and guards every route independently."),
        ("Hotspot Detection",
         "MongoDB's aggregation pipeline was used to group cases by (department × category) pairs and surface any pair with ≥5 open cases."),
        ("Working-Day Escalation",
         "The cron job needed to count business days (Mon–Fri) only. Implemented a helper that skips weekends when computing elapsed days."),
        ("Anonymous + Trackable",
         "Stripped PII from the stored complaint while issuing the submitter a unique Case ID stored locally, maintaining privacy and traceability."),
        ("CORS in Production",
         "Strict CORS policy required matching exact CLIENT_URL env variable in Express to prevent cross-origin blocking between Vercel and Render."),
    ]
    for i, (challenge, detail) in enumerate(items):
        y = 1.2 + i * 1.06
        add_rect(slide, 0.45, y+0.05, 2.5, 0.62, fill_color=RGBColor(0xFF, 0x5C, 0x5C))
        add_text(slide, challenge, 0.48, y+0.12, 2.44, 0.5,
                 font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, detail, 3.15, y+0.08, 9.8, 0.68, font_size=12, color=COLOR_LIGHT_GREY)


def slide_deployment(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=COLOR_CARD)
    add_text(slide, "Deployment Strategy", 0.6, 0.15, 12, 0.65,
             font_size=28, bold=True, color=COLOR_ACCENT)
    add_divider(slide, 1.05)

    add_text(slide, "GitHub Repository → Auto-trigger CI/CD on every push to main", 0.6, 1.15, 12, 0.5,
             font_size=14, color=COLOR_LIGHT_GREY, italic=True)

    boxes = [
        ("🌐 Vercel",         "Frontend Host",   ["Next.js frontend",  "Edge-optimised CDN", "Auto-deploy from GitHub", "NEXT_PUBLIC_API_BASE env var"]),
        ("🖥 Render",         "Backend Host",   ["Node.js Express API","Blueprint YAML config","Auto-restart on crash", "Environment secrets managed"]),
        ("🍃 MongoDB Atlas",  "Cloud Database", ["Cloud NoSQL cluster", "Secure connection string", "Auto-backup enabled", "Atlas IP white-listing"]),
    ]
    for i, (title, sub, items) in enumerate(boxes):
        x = 0.5 + i * 4.25
        add_rect(slide, x, 1.8, 3.9, 4.5, fill_color=COLOR_CARD,
                 line_color=COLOR_ACCENT, line_width=Pt(1.5))
        add_text(slide, title, x+0.15, 1.95, 3.6, 0.5,
                 font_size=17, bold=True, color=COLOR_ACCENT)
        add_text(slide, sub, x+0.15, 2.5, 3.6, 0.38,
                 font_size=12, color=COLOR_ACCENT2, italic=True)
        add_multiline(slide, items, x+0.15, 2.95, 3.6, 3.0,
                      font_size=12, leading_char="• ", leading_color=COLOR_ACCENT2)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=COLOR_CARD)
    add_text(slide, "Future Roadmap", 0.6, 0.15, 12, 0.65,
             font_size=28, bold=True, color=COLOR_ACCENT)
    add_divider(slide, 1.05)

    items = [
        ("🤖 AI Auto-Categorisation", "Use NLP to automatically tag complaints by department and urgency from free-text input."),
        ("🔔 Real-Time Notifications", "Add Socket.io for instant push alerts to Case Managers when a case is assigned or escalated."),
        ("📱 Mobile Companion App",    "React Native app using the existing REST API for on-the-go complaint submission and tracking."),
        ("🔗 Slack / Teams Integration","Webhook alerts to management channels when a hotspot crosses the threshold."),
    ]
    for i, (title, detail) in enumerate(items):
        y = 1.25 + i * 1.35
        add_rect(slide, 0.45, y, 12.4, 1.1, fill_color=COLOR_CARD,
                 line_color=COLOR_ACCENT2, line_width=Pt(1.2))
        add_text(slide, title,  0.7, y+0.08, 4.5, 0.5,
                 font_size=15, bold=True, color=COLOR_ACCENT)
        add_text(slide, detail, 5.3, y+0.1,  7.3, 0.8,
                 font_size=13, color=COLOR_LIGHT_GREY)

    add_text(slide, '"NeoConnect is the foundation. The roadmap is the skyline."',
             0.6, 6.6, 12, 0.55, font_size=13, color=COLOR_ACCENT2, italic=True)


# ─── BUILD ALL SLIDES ────────────────────────────────────────────────────────
slide_title(prs)
slide_problem(prs)
slide_solution(prs)
slide_approach(prs)
slide_techstack(prs)
slide_features(prs)
slide_challenges(prs)
slide_deployment(prs)
slide_roadmap(prs)

out_path = r"c:\Users\varsh\Downloads\Full Stack Developer - Use Case and submission\NeoConnect_Presentation.pptx"
prs.save(out_path)
print(f"✅ PPT saved to: {out_path}")
